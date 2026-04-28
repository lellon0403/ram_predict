"""
RAM Price Predictor - 발표용 PPT (React 테마 / 이미지 원본 비율 유지)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

IMG = os.path.join(os.path.dirname(__file__), 'image')

# ── 색상 ─────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E,0x40,0xAF)
BLUE_MID   = RGBColor(0x25,0x63,0xEB)
BLUE_LIGHT = RGBColor(0xEF,0xF6,0xFF)
BLUE_BDR   = RGBColor(0xBF,0xDB,0xFE)
ORANGE     = RGBColor(0xF9,0x73,0x16)
GREEN      = RGBColor(0x16,0xA3,0x4A)
GREEN_LT   = RGBColor(0xF0,0xFD,0xF4)
PURPLE     = RGBColor(0x6D,0x28,0xD9)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
GRAY       = RGBColor(0x47,0x55,0x69)
SLATE      = RGBColor(0x64,0x74,0x8B)
DARK       = RGBColor(0x1E,0x29,0x3B)
BG         = RGBColor(0xF0,0xF6,0xFF)
LIGHT      = RGBColor(0xF1,0xF5,0xF9)

TOTAL = 18   # 전체 슬라이드 수

# ── 유틸 ─────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(sl, x, y, w, h, fill=None, line=None):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def txt(sl, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def grid_bg(sl):
    rect(sl, 0, 0, 13.33, 7.5, fill=BG)

def header_bar(sl, title, sub=None):
    rect(sl, 0, 0, 13.33, 1.1, fill=WHITE)
    rect(sl, 0, 1.07, 13.33, 0.06, fill=BLUE_MID)
    rect(sl, 0.35, 0.18, 0.58, 0.58, fill=BLUE_MID)
    txt(sl, '💾', 0.35, 0.18, 0.58, 0.58, size=22, align=PP_ALIGN.CENTER)
    txt(sl, title, 1.05, 0.1, 11.0, 0.54, size=24, bold=True, color=DARK)
    if sub:
        txt(sl, sub, 1.05, 0.62, 11.0, 0.38, size=12, color=SLATE)

def page_num(sl, n):
    txt(sl, f'{n} / {TOTAL}', 12.5, 7.15, 0.8, 0.28,
        size=10, color=SLATE, align=PP_ALIGN.RIGHT)

def card(sl, x, y, w, h, fill=WHITE, border=BLUE_BDR):
    rect(sl, x, y, w, h, fill=fill, line=border)

def place_img(sl, fname, y_start=1.18, x_margin=0.22,
              max_w=12.89, max_h=6.1, caption=None):
    """원본 비율 유지, 슬라이드 내 중앙 배치"""
    path = os.path.join(IMG, fname)
    if not os.path.exists(path):
        return
    with PILImage.open(path) as im:
        ow, oh = im.size
    ratio = ow / oh
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    x = x_margin + (max_w - w) / 2
    # 캡션 공간 확보
    cap_h = 0.38 if caption else 0
    img_area_h = max_h - cap_h
    if h > img_area_h:
        h = img_area_h
        w = h * ratio
        x = x_margin + (max_w - w) / 2
    y = y_start + (img_area_h - h) / 2
    sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        txt(sl, caption, x_margin, y_start + max_h - cap_h + 0.02,
            max_w, cap_h - 0.05, size=10, color=SLATE, italic=True,
            align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 1 — 타이틀
# ════════════════════════════════════════════
sl = blank()
rect(sl, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
rect(sl, 9.2,  0, 4.13, 7.5, fill=BLUE_MID)
rect(sl, 10.5, 0, 2.83, 7.5, fill=RGBColor(0x1D,0x4E,0xD8))
txt(sl, '💾', 0.7, 0.9, 1.2, 1.0, size=48, color=WHITE)
txt(sl, 'RAM Price Predictor', 0.7, 1.85, 9.0, 0.9,
    size=40, bold=True, color=WHITE)
txt(sl, '삼성 DDR5 16GB  ·  AI 단일 변수 선형 회귀 예측 시스템',
    0.7, 2.75, 9.0, 0.55, size=18, color=RGBColor(0xBF,0xDB,0xFE))
for i, t in enumerate([
    '✅  Mission 1 — 데이터 수집 & Preprocessing',
    '✅  Mission 2 — TensorFlow Linear Regression',
    '✅  Mission 3 — n8n 서비스 완성',
]):
    txt(sl, t, 0.7, 3.55+i*0.52, 9.0, 0.45,
        size=15, color=RGBColor(0x93,0xC5,0xFD))
txt(sl, '2026. 04', 0.7, 5.6, 4.0, 0.4, size=13, color=SLATE)

# ════════════════════════════════════════════
# SLIDE 2 — 프로젝트 개요
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl, '프로젝트 개요', '삼성 DDR5 16GB RAM 가격 예측 서비스')
page_num(sl, 2)

txt(sl,'📌 배경 및 목적',0.5,1.22,12.5,0.4,size=15,bold=True,color=BLUE_MID)
rect(sl,0.5,1.62,12.3,0.04,fill=BLUE_BDR)
for i,t in enumerate([
    '🔹  AI 수요 급증으로 DDR5 16GB 가격이 2년 만에 약 4배 폭등 (6만원 → 33만원)',
    '🔹  소비자가 "지금 사는 게 맞는가?" 판단할 수 있는 AI 예측 서비스 개발',
    '🔹  180일 실제 다나와 시세 → TensorFlow 단일 변수 선형 회귀 모델 학습',
]):
    txt(sl,t,0.65,1.72+i*0.44,12.0,0.38,size=13,color=GRAY)

txt(sl,'📊 데이터 현황',0.5,3.12,12.5,0.4,size=15,bold=True,color=BLUE_MID)
rect(sl,0.5,3.52,12.3,0.04,fill=BLUE_BDR)
for i,(ic,title,val,sub,bg) in enumerate([
    ('📅','학습 데이터 기간','180일','2025-10-17 ~ 2026-04-14',BLUE_LIGHT),
    ('💰','가격 범위','95,285 ~ 330,461원','Samsung DDR5-5600 16GB',GREEN_LT),
    ('🏪','데이터 출처','다나와 최저가','일별 현금 최저가 기준',RGBColor(0xFA,0xF5,0xFF)),
    ('🔧','파이프라인','n8n → Flask → React','Webhook 기반 서비스',LIGHT),
]):
    x=0.5+i*3.1
    card(sl,x,3.62,2.95,2.1,fill=bg)
    txt(sl,ic+' '+title,x+0.15,3.72,2.7,0.38,size=11,bold=True,color=BLUE_DARK)
    txt(sl,val,x+0.15,4.18,2.7,0.52,size=16,bold=True,color=BLUE_MID)
    txt(sl,sub,x+0.15,4.78,2.7,0.38,size=10,color=GRAY)

# ════════════════════════════════════════════
# SLIDE 3 — 미션 달성
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'미션 달성 현황','3가지 요구사항 모두 충족')
page_num(sl,3)
for i,(num_lbl,title,detail,color) in enumerate([
    ('Mission 1','데이터 수집 & Preprocessing',
     'CSV 730일 수집 → 최근 180일 선별 → Z-score 정규화 → 날짜 인덱스 단일 변수화 → Train 80% / Val 20%', GREEN),
    ('Mission 2','TensorFlow Regression 모델',
     'Dense(1, activation=\'linear\') = H(x)=Wx+b  /  MSE 코스트 함수  /  Adam Optimizer  /  EarlyStopping', BLUE_MID),
    ('Mission 3','n8n 서비스 완성',
     'Webhook 노드 → HTTP Request(Flask /predict) → Respond to Webhook  /  React 프론트엔드 완성', ORANGE),
]):
    y=1.35+i*1.9
    card(sl,0.5,y,12.3,1.72)
    rect(sl,0.5,y,0.2,1.72,fill=color)
    txt(sl,'✅ '+num_lbl,0.85,y+0.12,2.8,0.42,size=13,bold=True,color=color)
    txt(sl,title,0.85,y+0.55,11.5,0.5,size=20,bold=True,color=DARK)
    txt(sl,detail,0.85,y+1.1,11.3,0.5,size=11,color=GRAY)

# ════════════════════════════════════════════
# SLIDE 4 — Mission 1: 전처리
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'Mission 1 — 데이터 전처리','Preprocessing Pipeline')
page_num(sl,4)
for i,(n,title,code,desc,color) in enumerate([
    ('①','데이터 수집 & 선별',
     'data/ram_prices.csv  730일 → 최근 180일',
     '최근 상승 구간이 선형에 더 가까워 단일 변수 회귀에 적합',BLUE_MID),
    ('②','Z-score 정규화',
     'x_sc=(x-x_mean)/x_std  /  y_sc=(y-y_mean)/y_std',
     '평균=0, 표준편차=1 변환 / 역변환 파라미터 scaler.json 저장',GREEN),
    ('③','단일 변수 X 생성',
     'X=[0,1,...,179]  →  y=[가격₀, 가격₁,...,가격₁₇₉]',
     'PDF의 H(x)=Wx+b 에 맞게 날짜 인덱스 1개 → 가격 1개',ORANGE),
    ('④','Train / Val 분리',
     'Train: 앞 80%(144일)  /  Val: 뒤 20%(36일)',
     '시계열 순서 유지 — 무작위 셔플 없이 앞→뒤 분리',PURPLE),
]):
    x=0.45+(i%2)*6.45; y=1.32+(i//2)*2.9
    card(sl,x,y,6.15,2.72)
    rect(sl,x,y,6.15,0.07,fill=color)
    txt(sl,n+'  '+title,x+0.18,y+0.15,5.85,0.42,size=15,bold=True,color=color)
    rect(sl,x+0.18,y+0.65,5.78,0.72,fill=DARK)
    txt(sl,code,x+0.32,y+0.73,5.5,0.6,size=11,color=RGBColor(0x93,0xC5,0xFD))
    txt(sl,desc,x+0.18,y+1.48,5.78,0.82,size=11,color=GRAY)

# ════════════════════════════════════════════
# SLIDE 5 — Mission 2: 모델 + 코스트 함수
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'Mission 2 — TensorFlow 모델 & 코스트 함수',
           'H(x)=W·x+b  /  MSE Cost Function  /  Adam Optimizer')
page_num(sl,5)

txt(sl,'모델 아키텍처',0.5,1.25,6.0,0.38,size=14,bold=True,color=BLUE_MID)
for i,(lbl,sub,bg,fc) in enumerate([
    ('Input  (shape: 1)','날짜 인덱스 1개',BLUE_LIGHT,BLUE_DARK),
    ('Dense(1, activation=\'linear\')','y = W·x + b',RGBColor(0xFE,0xF3,0xC7),RGBColor(0x92,0x40,0x0E)),
    ('Output  (shape: 1)','예측 가격 1개',GREEN_LT,GREEN),
]):
    y=1.7+i*1.22
    rect(sl,0.5,y,5.2,0.72,fill=bg,line=fc)
    txt(sl,lbl,0.65,y+0.08,4.9,0.38,size=13,bold=True,color=fc)
    txt(sl,sub,0.65,y+0.44,4.9,0.22,size=10,color=fc,italic=True)
    if i<2: txt(sl,'↓',2.55,y+0.78,0.5,0.35,size=18,color=SLATE,align=PP_ALIGN.CENTER)

txt(sl,'파라미터: W 1개 + b 1개 = 총 2개',0.5,5.42,6.0,0.35,size=12,bold=True,color=BLUE_DARK)
rect(sl,0.5,5.82,5.8,0.48,fill=DARK)
txt(sl,'model.compile(loss=\'mse\', optimizer=Adam(lr=0.001))',
    0.65,5.9,5.5,0.35,size=11,color=RGBColor(0x93,0xC5,0xFD))

# 오른쪽: 코스트 함수
card(sl,6.85,1.2,6.05,5.55)
txt(sl,'📐 Cost Function & Optimizer',7.05,1.3,5.7,0.42,size=14,bold=True,color=BLUE_MID)
rect(sl,7.05,1.72,5.7,0.04,fill=BLUE_BDR)

rect(sl,7.05,1.85,5.7,1.6,fill=BLUE_LIGHT)
txt(sl,'MSE Cost Function  (PDF와 동일)',7.22,1.93,5.42,0.38,size=12,bold=True,color=BLUE_DARK)
rect(sl,7.22,2.35,5.42,0.62,fill=DARK)
txt(sl,'cost(W,b) = (1/n) × Σ(H(xᵢ) - yᵢ)²',
    7.38,2.43,5.1,0.45,size=13,bold=True,color=RGBColor(0x93,0xC5,0xFD))
txt(sl,'예측값과 실제값 차이의 제곱 평균 → 최소화하도록 W, b 조정',
    7.22,3.02,5.42,0.35,size=10,color=GRAY)

rect(sl,7.05,3.58,5.7,1.6,fill=RGBColor(0xFE,0xF3,0xC7))
txt(sl,'Adam Optimizer  (Gradient Descent 개선)',7.22,3.66,5.42,0.38,size=12,bold=True,color=RGBColor(0x92,0x40,0x0E))
rect(sl,7.22,4.08,5.42,0.62,fill=DARK)
txt(sl,'W = W - α × ∂cost/∂W  (학습률 자동 조절)',
    7.38,4.16,5.1,0.45,size=12,color=RGBColor(0xFD,0xD8,0x8A))
txt(sl,'PDF SGD보다 학습률 자동 조절 → 더 빠르고 정밀하게 수렴',
    7.22,4.75,5.42,0.35,size=10,color=GRAY)

rect(sl,7.05,5.3,5.7,1.38,fill=LIGHT)
txt(sl,'하이퍼파라미터',7.22,5.38,5.42,0.35,size=12,bold=True,color=DARK)
for j,(k,v) in enumerate([('EPOCHS','2000  (EarlyStopping patience=100)'),
                           ('BATCH','32'),('LR','0.001 (Adam)')]):
    txt(sl,k+':',7.22,5.78+j*0.28,1.5,0.25,size=10,bold=True,color=BLUE_DARK)
    txt(sl,v,   8.75,5.78+j*0.28,3.95,0.25,size=10,color=GRAY)

# ════════════════════════════════════════════
# SLIDE 6 — 학습 결과
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'모델 학습 결과','Training & Validation Performance')
page_num(sl,6)
for i,(ic,title,val,sub,bg,fc) in enumerate([
    ('⏱','최적 에포크','222 / 2000','EarlyStopping 조기 종료',BLUE_LIGHT,BLUE_DARK),
    ('📉','검증 MAE','±10,821원','정규화 역변환 후 원화',GREEN_LT,GREEN),
    ('📊','학습 샘플','144개','180일 중 앞 80%',RGBColor(0xFE,0xF3,0xC7),RGBColor(0x92,0x40,0x0E)),
    ('🔍','검증 샘플','36개','180일 중 뒤 20%',RGBColor(0xFA,0xF5,0xFF),PURPLE),
]):
    x=0.5+i*3.1
    card(sl,x,1.28,2.92,1.88,fill=bg)
    txt(sl,ic,x+0.18,1.38,0.55,0.5,size=22,color=fc)
    txt(sl,title,x+0.75,1.43,2.0,0.38,size=11,bold=True,color=fc)
    txt(sl,val,x+0.18,1.9,2.62,0.55,size=22,bold=True,color=fc,align=PP_ALIGN.CENTER)
    txt(sl,sub,x+0.18,2.5,2.62,0.35,size=10,color=GRAY,align=PP_ALIGN.CENTER)

txt(sl,'📋 학습 과정 타임라인',0.5,3.32,12.5,0.38,size=14,bold=True,color=BLUE_MID)
rect(sl,0.5,3.7,12.3,0.04,fill=BLUE_BDR)
for i,(ep,desc,color) in enumerate([
    ('에포크 1~100','학습률 0.001로 초기 학습. val_loss 지속 감소.',GREEN),
    ('에포크 ~222','222 에포크에서 최적 가중치 달성 (best val_loss).',BLUE_MID),
    ('에포크 ~272','50회 미개선 → ReduceLROnPlateau: 학습률 절반 감소.',ORANGE),
    ('에포크 322','100회 추가 미개선 → EarlyStopping 발동, best weights 복원.',RGBColor(0xDC,0x26,0x26)),
]):
    y=3.82+i*0.72
    rect(sl,0.5,y,2.62,0.58,fill=LIGHT)
    rect(sl,0.5,y,0.12,0.58,fill=color)
    txt(sl,ep,0.72,y+0.1,2.38,0.38,size=11,bold=True,color=color)
    txt(sl,desc,3.28,y+0.12,9.5,0.36,size=12,color=GRAY)
    rect(sl,3.12,y+0.22,0.12,0.12,fill=color)

# ════════════════════════════════════════════
# SLIDE 7 — Mission 3: n8n
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'Mission 3 — n8n 서비스 구성','Webhook → Flask → Response 파이프라인')
page_num(sl,7)
for i,(ic,title,desc,color) in enumerate([
    ('🔗','Webhook 노드','POST /ram-predict\n외부 요청 수신',BLUE_MID),
    ('📡','HTTP Request 노드','POST 192.168.45.83:\n5000/predict 호출',ORANGE),
    ('📤','Respond to Webhook','JSON 응답 반환\nCORS 헤더 포함',GREEN),
]):
    x=0.9+i*4.0
    card(sl,x,1.32,3.6,2.35)
    rect(sl,x,1.32,3.6,0.08,fill=color)
    txt(sl,ic,x+0.2,1.48,0.65,0.6,size=26,color=color)
    txt(sl,title,x+0.9,1.52,2.55,0.6,size=14,bold=True,color=color)
    txt(sl,desc,x+0.2,2.18,3.2,0.88,size=11,color=GRAY)
    if i<2: txt(sl,'→',x+3.65,2.1,0.6,0.55,size=28,bold=True,color=SLATE)

rect(sl,0.5,3.88,12.3,0.04,fill=BLUE_BDR)
txt(sl,'⬇  전체 서비스 흐름',0.5,3.98,12.5,0.38,size=14,bold=True,color=BLUE_MID)
for i,(n,t) in enumerate([
    ('1️⃣','React 프론트엔드가 n8n Webhook URL로 POST 요청 전송'),
    ('2️⃣','n8n Webhook 수신 → HTTP Request 노드가 Flask /predict 호출'),
    ('3️⃣','Flask: 날짜 인덱스 생성 → 정규화 → TF 모델 추론 → 90일 예측'),
    ('4️⃣','n8n Respond 노드가 JSON 반환 → React 차트/카드 렌더링'),
]):
    txt(sl,n,0.5,4.45+i*0.52,0.45,0.45,size=14)
    txt(sl,t,1.05,4.45+i*0.52,11.7,0.45,size=13,color=GRAY)

# ════════════════════════════════════════════
# SLIDE 8 — 시연: 메인 화면
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'서비스 시연 ① — 메인 화면','n8n 연결 완료 후 표시 / 오늘 시세 ₩323,163')
page_num(sl,8)
place_img(sl,'main.png',
          caption='▲ 메인 화면: 오늘 시세 실시간 표시 · 구매 희망가 입력 · 현재가/-5%/-10%/-20% 퀵버튼')

# ════════════════════════════════════════════
# SLIDE 9 — 시연: 가격 입력
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'서비스 시연 ② — 가격 입력','구매 희망가 ₩330,000 입력 후 예측하기')
page_num(sl,9)
place_img(sl,'price_enter.png',
          caption='▲ 330,000원 입력 → 예측하기 버튼 활성화 · 실시간 원화 표기')

# ════════════════════════════════════════════
# SLIDE 10 — 시연: AI 구매 추천
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'서비스 시연 ③ — AI 구매 추천','목표가 이하 현재가 → 즉시 구매 권고')
page_num(sl,10)
place_img(sl,'recommand.png',
          caption='▲ 현재가(₩323,163) < 목표가(₩330,000) → "지금 사세요!" / 앞으로 오를 가능성 높음')

# ════════════════════════════════════════════
# SLIDE 11 — 시연: 기간별 예측 카드
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'서비스 시연 ④ — 기간별 예측 가격','1주일 / 1개월 / 3개월 예측')
page_num(sl,11)
place_img(sl,'ram_predict.png',
          caption='▲ 1주일 ₩345,837(+7%) · 1개월 ₩371,859(+15.1%) · 3개월 ₩439,740(+36.1%)')

# ════════════════════════════════════════════
# SLIDE 12 — 시연: 예측 차트
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'서비스 시연 ⑤ — 가격 추이 & 예측 차트','최근 90일 실제 + 90일 예측 직선')
page_num(sl,12)
place_img(sl,'ram_graph.png',
          caption='▲ 파란 실선=실제가 · 주황 점선=예측(직선) · 초록 점선=목표가 ₩330,000')

# ════════════════════════════════════════════
# SLIDE 13 — 예측 정확도: 모달 통계
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'예측 정확도 검증 ① — 통계 요약','헤더 🎯 버튼 클릭 시 모달 팝업')
page_num(sl,13)
place_img(sl,'actual_data_crossing.png',
          caption='▲ MAPE 8.53% / MAE ±27,185원 / 예측 방향 상승 ✓  (2026-04-15 ~ 04-27, 13일 비교)')

# ════════════════════════════════════════════
# SLIDE 14 — 예측 정확도: 비교 그래프
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'예측 정확도 검증 ② — 예측 vs 실제 그래프','두 선 모두 상승 추세 일치')
page_num(sl,14)
place_img(sl,'actual_graph.png',
          caption='▲ 주황 점선(예측)이 파란 실선(실제)보다 약 8% 높게 형성 — 계통 오차(Systematic Error)')

# ════════════════════════════════════════════
# SLIDE 15 — 예측 정확도: 날짜별 상세 표
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'예측 정확도 검증 ③ — 날짜별 상세 비교','오차율 -7.8% ~ -9.1% 범위에서 안정적으로 유지')
page_num(sl,15)
place_img(sl,'actual_matrix.png',
          caption='▲ 오차율이 일정(계통 오차) → 모델이 상승 추세를 과도하게 학습한 결과')

# ════════════════════════════════════════════
# SLIDE 16 — 오차율 분석 & 한계
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'오차율 8.53% 분석 — 원인 및 한계','단일 변수 선형 회귀의 구조적 특성')
page_num(sl,16)

# 왼쪽: 작동하는 것
txt(sl,'✅ MSE 코스트 함수 정상 작동',0.5,1.22,6.2,0.42,size=15,bold=True,color=GREEN)
rect(sl,0.5,1.65,6.0,0.04,fill=GREEN)
card(sl,0.5,1.78,6.0,1.35,fill=GREEN_LT)
txt(sl,'cost(W,b) = (1/n) × Σ(H(xᵢ) - yᵢ)²',
    0.65,1.87,5.75,0.45,size=14,bold=True,color=GREEN)
txt(sl,'Adam이 222 에포크 동안 이 값을 최소화\n검증 MAE ±10,821원 달성 → MSE는 정상 동작',
    0.65,2.38,5.75,0.65,size=11,color=GRAY)

# 왼쪽: 오차 원인
txt(sl,'❌ 오차가 남는 이유',0.5,3.28,6.2,0.42,size=15,bold=True,color=RGBColor(0xDC,0x26,0x26))
rect(sl,0.5,3.7,6.0,0.04,fill=RGBColor(0xDC,0x26,0x26))
for i,(title,desc) in enumerate([
    ('① 구조적 한계','H(x)=Wx+b 는 직선만 표현 가능\n가격 곡선 패턴은 직선으로 완벽히 맞출 수 없음'),
    ('② 계통 오차','예측이 일관되게 +8% 높음\n상승 구간(180일)을 과도하게 학습'),
    ('③ 단일 변수','날짜 인덱스 하나만 입력\n요일·거래량·계절성 등 다른 영향 변수 무시'),
]):
    y=3.85+i*1.02
    card(sl,0.5,y,6.0,0.92,fill=RGBColor(0xFE,0xF2,0xF2))
    txt(sl,title,0.68,y+0.1,2.5,0.35,size=12,bold=True,color=RGBColor(0xDC,0x26,0x26))
    txt(sl,desc,0.68,y+0.46,5.65,0.4,size=10,color=GRAY)

# 오른쪽: 시각 다이어그램
card(sl,6.85,1.22,6.05,3.58)
txt(sl,'📊 직선 vs 실제 가격 곡선',7.05,1.32,5.7,0.38,size=14,bold=True,color=BLUE_MID)
rect(sl,7.05,1.7,5.7,0.04,fill=BLUE_BDR)
rect(sl,7.05,1.82,5.7,2.88,fill=DARK)
for i,(label,val,color) in enumerate([
    ('실제:','   ╭──────── ← 33만원',RGBColor(0x93,0xC5,0xFD)),
    ('','  ╱  ← 급등',              RGBColor(0x93,0xC5,0xFD)),
    ('예측:','─────────── ← 예측직선',ORANGE),
    ('','   ↑ 이 격차 = 8% 오차',   RGBColor(0xFE,0xF0,0x8A)),
    ('','H(x)=Wx+b = 직선만 가능',   WHITE),
]):
    txt(sl,label,7.2,  2.0+i*0.5,1.6,0.45,size=11,bold=bool(label),color=color)
    txt(sl,val,  8.85, 2.0+i*0.5,2.8,0.45,size=11,color=color)

# 오른쪽: 개선방향
card(sl,6.85,4.95,6.05,1.78)
txt(sl,'💡 개선 방향 (참고)',7.05,5.05,5.7,0.38,size=13,bold=True,color=ORANGE)
for i,t in enumerate([
    '다변수 회귀: 요일·거래량·계절성 추가 입력',
    '다항 회귀: H(x)=W₂x²+W₁x+b 곡선 피팅',
    '※ 본 과제는 단일 변수 선형 회귀 구조 준수',
]):
    txt(sl,('• ' if i<2 else '※ ')+t,
        7.05,5.48+i*0.42,5.9,0.38,
        size=11,color=(GRAY if i<2 else ORANGE),bold=(i==2))

# ════════════════════════════════════════════
# SLIDE 17 — 전체 아키텍처
# ════════════════════════════════════════════
sl = blank(); grid_bg(sl)
header_bar(sl,'전체 시스템 아키텍처','Data → Model → Service → UI')
page_num(sl,17)
for i,(title,detail,bg,fc) in enumerate([
    ('📦 데이터 레이어','data/ram_prices.csv  (최근 180일 다나와 시세)',BLUE_LIGHT,BLUE_DARK),
    ('🧠 모델 레이어','model/trainer.py  →  saved_model/model.keras  (TensorFlow Keras, 파라미터 2개)',RGBColor(0xFE,0xF3,0xC7),RGBColor(0x92,0x40,0x0E)),
    ('⚙️ 서버 레이어','app.py  Flask REST API  (POST /predict  /  GET /comparison)',GREEN_LT,GREEN),
    ('🔗 오케스트레이션','n8n  Webhook → HTTP Request → Respond  3-노드 파이프라인',RGBColor(0xFA,0xF5,0xFF),PURPLE),
    ('🖥️ UI 레이어','React + Vite + Recharts  (localhost:3001)  /  예측 정확도 모달 포함',LIGHT,GRAY),
]):
    y=1.32+i*1.06
    card(sl,0.5,y,12.3,0.92,fill=bg)
    rect(sl,0.5,y,0.18,0.92,fill=fc)
    txt(sl,title,0.82,y+0.08,3.5,0.36,size=13,bold=True,color=fc)
    txt(sl,detail,4.5,y+0.1,8.1,0.36,size=12,color=GRAY)
    if i<4: txt(sl,'↓',6.5,y+0.93,0.5,0.18,size=12,color=SLATE,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# SLIDE 18 — 결론
# ════════════════════════════════════════════
sl = blank()
rect(sl,0,0,13.33,7.5,fill=BLUE_DARK)
rect(sl,0,5.5,13.33,2.0,fill=BLUE_MID)
txt(sl,'결론 및 성과',0.7,0.5,12.0,0.75,size=34,bold=True,color=WHITE)
rect(sl,0.7,1.28,3.5,0.07,fill=ORANGE)
for i,(ic,title,desc) in enumerate([
    ('✅','Mission 1','180일 다나와 시세 / Z-score 정규화 / 단일 변수(날짜 인덱스) 전처리'),
    ('✅','Mission 2','Dense(1,linear) / MSE 코스트 함수 / Adam Optimizer / 파라미터 2개'),
    ('✅','Mission 3','n8n Webhook→Flask→React 3-tier 서비스 파이프라인 완성'),
    ('🎯','예측 정확도','MAPE 8.53% / MAE ±27,185원 / 상승 방향 정확 (13일 실제 비교)'),
    ('💡','오차 원인','H(x)=Wx+b 직선으로 곡선 패턴 표현 불가 — 단일 변수 선형 회귀의 구조적 한계'),
]):
    y=1.5+i*0.88
    txt(sl,ic,  0.7,y,0.5, 0.75,size=20,color=WHITE)
    txt(sl,title,1.2,y,2.2, 0.38,size=13,bold=True,color=ORANGE)
    txt(sl,desc, 1.2,y+0.42,11.2,0.38,size=12,color=RGBColor(0xBF,0xDB,0xFE))
txt(sl,'감사합니다',0.7,5.82,11.9,0.7,
    size=28,bold=True,color=WHITE,align=PP_ALIGN.CENTER)

# ── 저장 ────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(__file__), 'RAM_Price_Predictor_발표_v3.pptx')
prs.save(OUT)
print(f'저장 완료: {OUT}  ({len(prs.slides)} 슬라이드)')
